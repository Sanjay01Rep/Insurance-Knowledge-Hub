"""Extract, chunk, and index uploaded documents."""

from __future__ import annotations

import re
import threading
from pathlib import Path
from typing import Any, Callable

from db import logger

PROCESS_FAIL = "We couldn't process this document. Please check the file and try again."
NO_TEXT = "No readable text was found in this file."

TARGET_CHARS = 900
OVERLAP_CHARS = 150
EXCEL_ROWS_PER_CHUNK = 10

_lock = threading.Lock()
_in_flight: set[str] = set()


def enqueue_processing(doc_id: str) -> None:
    with _lock:
        if doc_id in _in_flight:
            return
        _in_flight.add(doc_id)
    thread = threading.Thread(target=_run_safe, args=(doc_id,), daemon=True)
    thread.start()


def enqueue_unfinished() -> None:
    from db import list_unfinished_documents

    for doc in list_unfinished_documents():
        enqueue_processing(doc["id"])


def _run_safe(doc_id: str) -> None:
    try:
        process_document(doc_id)
    except Exception:
        from db import update_document_status

        logger.exception("processing crashed id=%s", doc_id)
        update_document_status(doc_id, "failed", chunk_count=0, error_summary=PROCESS_FAIL)
    finally:
        with _lock:
            _in_flight.discard(doc_id)


def process_document(doc_id: str) -> None:
    from db import get_document, update_document_status
    from indexer import index_chunks

    doc = get_document(doc_id)
    if not doc:
        logger.info("processing skipped, missing id=%s", doc_id)
        return
    update_document_status(doc_id, "processing", error_summary="")
    logger.info("processing started id=%s type=%s", doc_id, doc.get("file_type"))

    path = Path(doc["storage_path"])
    if not path.exists():
        logger.info("processing failed, file missing id=%s", doc_id)
        update_document_status(doc_id, "failed", chunk_count=0, error_summary=PROCESS_FAIL)
        return

    file_type = (doc.get("file_type") or path.suffix.lstrip(".")).lower()
    extractors: dict[str, Callable[[Path, dict], list[dict[str, Any]]]] = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
        "txt": _extract_txt,
    }
    extractor = extractors.get(file_type)
    if extractor is None:
        update_document_status(doc_id, "failed", chunk_count=0, error_summary=PROCESS_FAIL)
        return

    try:
        chunks = extractor(path, doc)
    except Exception:
        logger.exception("extract failed id=%s type=%s", doc_id, file_type)
        update_document_status(doc_id, "failed", chunk_count=0, error_summary=PROCESS_FAIL)
        return

    if not chunks:
        logger.info("processing failed, no chunks id=%s", doc_id)
        update_document_status(doc_id, "failed", chunk_count=0, error_summary=NO_TEXT)
        return

    try:
        count = index_chunks(doc_id, chunks)
    except Exception:
        logger.exception("index failed id=%s", doc_id)
        update_document_status(doc_id, "failed", chunk_count=0, error_summary=PROCESS_FAIL)
        return

    update_document_status(doc_id, "ready", chunk_count=count, error_summary="")
    logger.info("processing ready id=%s chunks=%s", doc_id, count)


def _base_meta(doc: dict, extra: dict[str, Any] | None = None) -> dict[str, Any]:
    meta = {
        "document_id": doc["id"],
        "document_name": doc["name"],
        "file_type": doc.get("file_type") or "",
        "module": doc.get("module") or "",
        "topic": doc.get("topic") or "",
        "sheet_name": "",
        "headers": "",
        "source_locator": "",
        "chunk_index": 0,
    }
    if extra:
        meta.update(extra)
    return meta


def _make_chunk(doc: dict, index: int, content: str, extra: dict[str, Any]) -> dict[str, Any]:
    meta = _base_meta(doc, extra)
    meta["chunk_index"] = index
    return {"chunk_index": index, "content": content.strip(), "metadata": meta}


# --- TXT / generic prose -------------------------------------------------

def _extract_txt(path: Path, doc: dict) -> list[dict[str, Any]]:
    raw = path.read_bytes()
    text = None
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if not text:
        return []
    return _chunk_prose(doc, text, locator_prefix="Notes")


def _chunk_prose(doc: dict, text: str, *, locator_prefix: str, section: str = "") -> list[dict[str, Any]]:
    blocks = _split_keep_context(text)
    chunks: list[dict[str, Any]] = []
    for block in blocks:
        locator = locator_prefix
        if section:
            locator = f"{locator_prefix} · {section}"
        extra = {"source_locator": locator}
        if section:
            extra["section"] = section
        chunks.append(_make_chunk(doc, len(chunks), block, extra))
    return chunks


def _split_keep_context(text: str) -> list[str]:
    """Split on paragraphs, then sentences. Keep overlap. Avoid mid-sentence cuts."""
    cleaned = re.sub(r"\r\n?", "\n", text).strip()
    if not cleaned:
        return []
    paragraphs = [part.strip() for part in re.split(r"\n\s*\n", cleaned) if part.strip()]
    units: list[str] = []
    for para in paragraphs:
        if len(para) <= TARGET_CHARS:
            units.append(para)
        else:
            units.extend(_split_sentences(para))
    if not units:
        return [cleaned[:TARGET_CHARS]]

    chunks: list[str] = []
    current = ""
    for unit in units:
        candidate = f"{current}\n\n{unit}".strip() if current else unit
        if len(candidate) <= TARGET_CHARS or not current:
            current = candidate
            continue
        chunks.append(current)
        overlap = current[-OVERLAP_CHARS:] if len(current) > OVERLAP_CHARS else current
        current = f"{overlap}\n\n{unit}".strip()
    if current:
        chunks.append(current)
    return chunks


def _split_sentences(text: str) -> list[str]:
    parts = re.split(r"(?<=[.!?])\s+(?=[A-Z0-9])", text)
    return [part.strip() for part in parts if part.strip()]


# --- PDF -----------------------------------------------------------------

def _extract_pdf(path: Path, doc: dict) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    try:
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                locator = f"Page {page_number}"
                tables = page.extract_tables() or []
                for table in tables:
                    table_text = _table_to_text(table)
                    if table_text:
                        chunks.append(
                            _make_chunk(
                                doc,
                                len(chunks),
                                table_text,
                                {"source_locator": f"{locator} · Table"},
                            )
                        )
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    for block in _split_keep_context(page_text):
                        chunks.append(
                            _make_chunk(
                                doc,
                                len(chunks),
                                block,
                                {"source_locator": locator},
                            )
                        )
        if chunks:
            return chunks
    except Exception:
        logger.exception("pdfplumber failed, trying pypdf name=%s", doc.get("name"))

    from pypdf import PdfReader

    reader = PdfReader(str(path))
    for page_number, page in enumerate(reader.pages, start=1):
        page_text = (page.extract_text() or "").strip()
        if not page_text:
            continue
        for block in _split_keep_context(page_text):
            chunks.append(
                _make_chunk(
                    doc,
                    len(chunks),
                    block,
                    {"source_locator": f"Page {page_number}"},
                )
            )
    return chunks


def _table_to_text(table: list[list[Any]]) -> str:
    rows = [[("" if cell is None else str(cell).strip()) for cell in row] for row in table]
    rows = [row for row in rows if any(cell for cell in row)]
    if not rows:
        return ""
    headers = [h or f"Column {i + 1}" for i, h in enumerate(rows[0])]
    header_line = " | ".join(headers)
    body = []
    for row in rows[1:]:
        pairs = [f"{headers[i]}: {row[i]}" for i in range(min(len(headers), len(row))) if row[i]]
        if pairs:
            body.append(" | ".join(pairs))
    if not body:
        return f"Columns: {header_line}"
    return f"Columns: {header_line}\n" + "\n".join(body)


# --- DOCX ----------------------------------------------------------------

def _extract_docx(path: Path, doc: dict) -> list[dict[str, Any]]:
    from docx import Document

    document = Document(str(path))
    chunks: list[dict[str, Any]] = []
    current_heading = ""
    buffer: list[str] = []

    def flush() -> None:
        nonlocal buffer
        text = "\n".join(buffer).strip()
        buffer = []
        if not text:
            return
        section = current_heading or "Document"
        for block in _split_keep_context(text):
            chunks.append(
                _make_chunk(
                    doc,
                    len(chunks),
                    block,
                    {
                        "source_locator": f"Section: {section}",
                        "section": section,
                    },
                )
            )

    for para in document.paragraphs:
        style = (para.style.name if para.style else "") or ""
        text = (para.text or "").strip()
        if not text:
            continue
        if style.lower().startswith("heading"):
            flush()
            current_heading = text
            continue
        buffer.append(text)
    flush()

    for table in document.tables:
        grid = []
        for row in table.rows:
            grid.append([cell.text.strip() for cell in row.cells])
        table_text = _table_to_text(grid)
        if table_text:
            section = current_heading or "Table"
            chunks.append(
                _make_chunk(
                    doc,
                    len(chunks),
                    table_text,
                    {"source_locator": f"Section: {section} · Table", "section": section},
                )
            )
    return chunks


# --- PPTX ----------------------------------------------------------------

def _extract_pptx(path: Path, doc: dict) -> list[dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    chunks: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
            texts.append(title)
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            value = (shape.text or "").strip()
            if value and value not in texts:
                texts.append(value)
        body = "\n".join(texts).strip()
        if not body:
            continue
        locator = f"Slide {index}"
        if title:
            locator = f"Slide {index} · {title}"
        for block in _split_keep_context(body):
            chunks.append(
                _make_chunk(
                    doc,
                    len(chunks),
                    block,
                    {"source_locator": locator, "section": title},
                )
            )
    return chunks


# --- Excel (highest priority): workbook → sheet → headers → rows ---------

def _extract_xlsx(path: Path, doc: dict) -> list[dict[str, Any]]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=True, read_only=True)
    chunks: list[dict[str, Any]] = []
    workbook_name = doc["name"]
    try:
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [_cell_str(value) for value in row]
                if not any(values):
                    continue
                rows.append(values)
            if not rows:
                continue

            header_idx = 0
            headers = [_header_name(value, i) for i, value in enumerate(rows[header_idx])]
            data_rows = rows[header_idx + 1 :]
            if not data_rows:
                content = (
                    f"Workbook: {workbook_name}\n"
                    f"Sheet: {sheet_name}\n"
                    f"Columns: {' | '.join(headers)}"
                )
                chunks.append(
                    _make_chunk(
                        doc,
                        len(chunks),
                        content,
                        {
                            "sheet_name": sheet_name,
                            "headers": " | ".join(headers),
                            "source_locator": f"Sheet: {sheet_name}",
                        },
                    )
                )
                continue

            for start in range(0, len(data_rows), EXCEL_ROWS_PER_CHUNK):
                batch = data_rows[start : start + EXCEL_ROWS_PER_CHUNK]
                excel_row_start = start + 2
                excel_row_end = start + 1 + len(batch)
                lines = [
                    f"Workbook: {workbook_name}",
                    f"Sheet: {sheet_name}",
                    f"Columns: {' | '.join(headers)}",
                ]
                for offset, row in enumerate(batch):
                    pairs = []
                    for i, header in enumerate(headers):
                        cell = row[i] if i < len(row) else ""
                        if cell:
                            pairs.append(f"{header}: {cell}")
                    if pairs:
                        lines.append(f"Row {excel_row_start + offset}: " + " | ".join(pairs))
                content = "\n".join(lines)
                if content.count("\n") < 3:
                    continue
                chunks.append(
                    _make_chunk(
                        doc,
                        len(chunks),
                        content,
                        {
                            "sheet_name": sheet_name,
                            "headers": " | ".join(headers),
                            "source_locator": f"Sheet: {sheet_name} · rows {excel_row_start}-{excel_row_end}",
                        },
                    )
                )
        return chunks
    finally:
        workbook.close()


def _cell_str(value: Any) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _header_name(value: str, index: int) -> str:
    return value if value else f"Column {index + 1}"
